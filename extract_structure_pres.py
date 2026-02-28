import json
import re
from pptx import Presentation

PLACEHOLDER_PATTERN = re.compile(r"\{\{[^}]+\}\}")
TITLE_TYPE_PATTERN = re.compile(r"\{\{TITLE_([^}]+)\}\}")

def extract_placeholders_from_text(text):
    return PLACEHOLDER_PATTERN.findall(text)


def extract_slide_type(placeholders):
    """
    Ищем {{TITLE_XXXX}} и возвращаем XXXX в lower-case.
    Если не найдено — возвращаем None.
    """
    for ph in placeholders:
        match = TITLE_TYPE_PATTERN.match(ph)
        if match:
            return match.group(1)
    return None


def generate_structure_from_pptx(pptx_path):
    prs = Presentation(pptx_path)

    slides_structure = []

    for slide_index, slide in enumerate(prs.slides):
        replacements = {}
        all_placeholders = []

        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue

            text = shape.text
            placeholders = extract_placeholders_from_text(text)

            for ph in placeholders:
                all_placeholders.append(ph)
                if ph not in replacements:
                    replacements[ph] = ""

        if replacements:
            slide_type = extract_slide_type(all_placeholders)

            slides_structure.append({
                "slide_index": slide_index,
                "slide_type": slide_type,
                "replacements": replacements
            })

    return {"slides": slides_structure}


if __name__ == "__main__":
    input_pptx = "test.pptx"
    output_json = "structure.json"

    structure = generate_structure_from_pptx(input_pptx)

    with open(output_json, "w", encoding="utf-8") as f:
        json.dump(structure, f, ensure_ascii=False, indent=2)

    print(f"Структура сохранена в {output_json}")